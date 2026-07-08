from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

INDIGO = RGBColor(0x4F, 0x46, 0xE5)
DARK = RGBColor(0x1F, 0x29, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF3, 0xF4, 0xF6)
GRAY = RGBColor(0x6B, 0x72, 0x80)
EMERALD = RGBColor(0x05, 0x9C, 0x69)
ROSE = RGBColor(0xE0, 0x3E, 0x3E)
AMBER = RGBColor(0xD9, 0x77, 0x06)
SLATE = RGBColor(0x47, 0x55, 0x69)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_textbox(slide, left, top, width, height, items, font_size=16, color=DARK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_section_title(slide, title):
    add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), INDIGO)
    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8), title, font_size=32, bold=True, color=WHITE)

def add_accent_line(slide, left, top, width):
    shape = add_shape_bg(slide, left, top, width, Inches(0.04), INDIGO)

# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, INDIGO)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.3), prs.slide_height, RGBColor(0x37, 0x2C, 0xC8))
add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2), 'CovidCare', font_size=56, bold=True, color=WHITE)
add_textbox(slide, Inches(1.5), Inches(2.7), Inches(10), Inches(0.8), 'Vaccination Management System', font_size=36, bold=False, color=RGBColor(0xC7, 0xD2, 0xFE))
add_shape_bg(slide, Inches(1.5), Inches(3.7), Inches(2), Inches(0.06), WHITE)
add_textbox(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(0.6), 'eProject Report', font_size=24, bold=False, color=RGBColor(0xA5, 0xB4, 0xFC))
add_textbox(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.5), 'June 2026 - July 2026', font_size=18, color=RGBColor(0xC7, 0xD2, 0xFE))
add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.5), 'Developed by Team CovidCare', font_size=18, color=RGBColor(0xC7, 0xD2, 0xFE))

# ============================================================
# SLIDE 2: TEAM MEMBERS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Team Members')

members = [
    ('Muhammad Owais', 'Lead Developer & Architect', '(Super Admin)'),
    ('Muhammad Hammad', 'Full Stack Developer'),
    ('Muhammad Hunain', 'Frontend Developer'),
    ('Muhammad Hamayl', 'Backend Developer'),
    ('Muhammad Ali', 'UI/UX Designer'),
]
y = Inches(1.6)
for i, m in enumerate(members):
    color = INDIGO if i == 0 else DARK
    add_shape_bg(slide, Inches(1), y, Inches(11), Inches(0.9), LIGHT_BG)
    # accent bar
    add_shape_bg(slide, Inches(1), y, Inches(0.08), Inches(0.9), INDIGO)
    add_textbox(slide, Inches(1.5), y + Inches(0.1), Inches(4), Inches(0.5), m[0], font_size=20, bold=True, color=DARK)
    add_textbox(slide, Inches(5.5), y + Inches(0.1), Inches(5), Inches(0.5), m[1], font_size=18, color=GRAY)
    if len(m) > 2:
        add_textbox(slide, Inches(10), y + Inches(0.1), Inches(2), Inches(0.5), m[2], font_size=14, bold=True, color=INDIGO)
    y += Inches(1.05)

# ============================================================
# SLIDE 3: PROJECT SYNOPSIS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Project Synopsis')

synopsis = [
    'CovidCare is a web-based Vaccination Management System built with Laravel PHP framework.',
    'It serves three distinct user roles: Patients, Hospitals, and Administrators.',
    'Patients can register, search for hospitals, book test/vaccination appointments, and track their records.',
    'Hospitals can manage appointment requests, update test results, and record vaccination status.',
    'Administrators oversee the entire system - verifying hospitals, managing vaccine inventory, and generating reports.',
    'The system features role-based authentication, real-time notifications, and a responsive UI.',
    'Tech Stack: Laravel 11, PHP 8.x, MySQL, Tailwind CSS, Alpine.js, Vite, Bootstrap 5.',
]
add_bullet_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.5), synopsis, font_size=18, color=DARK)

# ============================================================
# SLIDE 4: SYSTEM ANALYSIS - REQUIREMENTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'System Analysis - Requirements')

add_textbox(slide, Inches(1), Inches(1.6), Inches(5), Inches(0.5), 'Functional Requirements', font_size=22, bold=True, color=INDIGO)
func_reqs = [
    'Multi-role user authentication (Patient, Hospital, Admin)',
    'Patient registration & profile management',
    'Hospital registration with admin approval workflow',
    'Hospital search by name/location',
    'Appointment booking (Test / Vaccination)',
    'Test result update by hospitals',
    'Vaccination status tracking',
    'Admin dashboard with statistics',
    'Vaccine inventory management',
    'Report generation & Excel export',
    'Notification system for all roles',
]
add_bullet_textbox(slide, Inches(1), Inches(2.2), Inches(5.5), Inches(5), func_reqs, font_size=14, color=DARK)

add_textbox(slide, Inches(7), Inches(1.6), Inches(5), Inches(0.5), 'Non-Functional Requirements', font_size=22, bold=True, color=INDIGO)
nonfunc_reqs = [
    'Responsive design for mobile & desktop',
    'Secure password hashing (bcrypt)',
    'Role-based access control',
    'CSRF & XSS protection',
    'Database indexing for performance',
    'Session management with guards',
    'Queue support for notifications',
    'SQLite in-memory testing support',
]
add_bullet_textbox(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(5), nonfunc_reqs, font_size=14, color=DARK)

# ============================================================
# SLIDE 5: SYSTEM ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'System Architecture')

arch_items = [
    'Presentation Layer: Blade Templates + Tailwind CSS + Alpine.js',
    '  - Responsive views for each role (Patient, Hospital, Admin)',
    '  - Landing page with theme customization',
    '',
    'Application Layer: Laravel MVC Framework',
    '  - Routes: 37 named routes across 3 role prefixes',
    '  - Controllers: 6 main controllers (3 Auth + 3 Feature)',
    '  - Middleware: auth guards for role-based access',
    '',
    'Data Layer: Eloquent ORM + MySQL',
    '  - 5 Models: User, Hospital, Vaccine, Appointment, Notification',
    '  - 8 Migration files defining the complete schema',
    '  - 2 Auth guards: web (patients/admins), hospital',
    '',
    'Infrastructure:',
    '  - Vite for asset bundling & hot reload',
    '  - Queue system for async notifications',
    '  - Session-based authentication',
]
add_bullet_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.5), arch_items, font_size=16, color=DARK)

# ============================================================
# SLIDE 6: DATABASE DESIGN - ER DIAGRAM (Text)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Database Design - Entity Relationship')

# Draw tables as shapes
tables_data = [
    ('users', ['id (PK)', 'name', 'email', 'password', 'role', 'phone', 'address', 'location'], Inches(0.8), Inches(1.6)),
    ('hospitals', ['id (PK)', 'hospital_name', 'email', 'password', 'address', 'location', 'status'], Inches(4.8), Inches(1.6)),
    ('vaccines', ['id (PK)', 'vaccine_name', 'status'], Inches(8.8), Inches(1.6)),
    ('appointments', ['id (PK)', 'patient_id (FK)', 'hospital_id (FK)', 'type', 'appointment_date', 'status', 'test_result', 'vaccination_status'], Inches(3.8), Inches(4.5)),
    ('notifications', ['id (PK)', 'notifiable_type', 'notifiable_id', 'type', 'title', 'message', 'link', 'is_read'], Inches(0.8), Inches(4.5)),
]

colors = [INDIGO, EMERALD, AMBER, ROSE, SLATE]

for (tname, cols, left, top), color in zip(tables_data, colors):
    h = Inches(0.4 + len(cols) * 0.35)
    shape = add_shape_bg(slide, left, top, Inches(3.6), Inches(0.4), color)
    add_textbox(slide, left, top + Inches(0.02), Inches(3.6), Inches(0.36), tname, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    for j, col in enumerate(cols):
        y = top + Inches(0.4 + j * 0.35)
        bg_c = LIGHT_BG if j % 2 == 0 else WHITE
        add_shape_bg(slide, left, y, Inches(3.6), Inches(0.35), bg_c)
        fc = EMERALD if 'PK' in col else (AMBER if 'FK' in col else DARK)
        add_textbox(slide, left + Inches(0.15), y + Inches(0.03), Inches(3.3), Inches(0.3), col, font_size=11, color=fc)

# Draw relationship lines (simple indicators)
# users 1---* appointments
add_textbox(slide, Inches(2.6), Inches(4.2), Inches(1.5), Inches(0.4), '1 ---< *', font_size=14, bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(6.8), Inches(4.2), Inches(1.5), Inches(0.4), '* >--- 1', font_size=14, bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: DATABASE RELATIONSHIPS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Database Relationships & Schema Summary')

rels = [
    'users (1) ----< appointments >---- (1) hospitals',
    '  - appointments.patient_id -> users.id (CASCADE on delete)',
    '  - appointments.hospital_id -> hospitals.id (CASCADE on delete)',
    '',
    'notifications (Polymorphic Relationship)',
    '  - notifiable_type = "App\\Models\\User" -> users table',
    '  - notifiable_type = "App\\Models\\Hospital" -> hospitals table',
    '',
    'Key Design Decisions:',
    '  - Single users table with "role" column for patient/admin distinction',
    '  - Separate hospitals table (not in users) due to different auth guard & fields',
    '  - Polymorphic notifications to support both User and Hospital recipients',
    '  - Appointments as central transactional hub connecting patients & hospitals',
    '  - Vaccines as standalone table for inventory management',
]
add_bullet_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.5), rels, font_size=16, color=DARK)

# ============================================================
# SLIDE 8: DATA FLOW DIAGRAM (Level 0)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Data Flow Diagram - Level 0 (Context Diagram)')

# Draw entities
entities = [
    ('Patient', Inches(1), Inches(3)),
    ('Hospital', Inches(10.5), Inches(1.5)),
    ('Admin', Inches(10.5), Inches(5)),
]
for name, left, top in entities:
    shape = add_shape_bg(slide, left, top, Inches(2), Inches(1), INDIGO)
    add_textbox(slide, left, top + Inches(0.2), Inches(2), Inches(0.6), name, font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Center system box
sys_box = add_shape_bg(slide, Inches(4.5), Inches(2.5), Inches(4.5), Inches(2.5), LIGHT_BG)
sys_box.line.color.rgb = INDIGO
sys_box.line.width = Pt(2)
add_textbox(slide, Inches(4.7), Inches(3), Inches(4.1), Inches(0.6), 'CovidCare System', font_size=22, bold=True, color=INDIGO, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.7), Inches(3.7), Inches(4.1), Inches(0.8), 'Registration\nAppointment Booking\nResults & Reports', font_size=14, color=DARK, alignment=PP_ALIGN.CENTER)

# Flow labels
flows = [
    (Inches(3), Inches(3.3), 'Registers, Books,\nViews Status'),
    (Inches(9), Inches(1.7), 'Manages Requests,\nUpdates Results'),
    (Inches(9), Inches(5.3), 'Verifies, Manages\nVaccines, Reports'),
]
for left, top, text in flows:
    add_textbox(slide, left, top, Inches(1.8), Inches(0.8), text, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5), 'External entities interact with the system through their respective auth guards and feature controllers.', font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: USE CASE DIAGRAM (Text)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Use Case Diagram - Actors & Features')

use_cases = [
    'Patient (Actor):',
    '  - Register Account',
    '  - Login / Logout',
    '  - Search Hospitals by name/location',
    '  - Book Appointments (Test / Vaccination)',
    '  - View Appointment History',
    '  - Manage Profile (update/delete account)',
    '  - Receive Notifications',
    '',
    'Hospital (Actor):',
    '  - Register Account (status: pending)',
    '  - Login / Logout (only when approved)',
    '  - View Appointment Requests',
    '  - Approve / Reject Appointments',
    '  - Update Test Results (Negative/Positive)',
    '  - Update Vaccination Status',
    '  - Receive Notifications',
    '',
    'Admin (Actor):',
    '  - Login / Logout',
    '  - View Dashboard with system stats',
    '  - Verify / Approve / Reject Hospitals',
    '  - Manage Vaccines (add, toggle availability)',
    '  - View All Appointments Reports',
    '  - Export Data to Excel',
]
add_bullet_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.5), use_cases, font_size=14, color=DARK)

# ============================================================
# SLIDE 10: APPLICATION FLOW - PATIENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Application Flow - Patient Module')

flow_steps = [
    'Step 1: Register / Login',
    '  - New users register with name, email, password, phone, address',
    '  - Existing users login with email + password (web guard, role=patient)',
    '',
    'Step 2: Search for Hospitals',
    '  - Navigate to Search page',
    '  - Search by hospital name or location',
    '  - System returns list of approved hospitals',
    '',
    'Step 3: Book Appointment',
    '  - Select a hospital from search results',
    '  - Choose appointment type: COVID Test or Vaccination',
    '  - Pick a date and submit',
    '  - Hospital receives notification of new booking',
    '',
    'Step 4: Track Appointments',
    '  - View all appointments on Appointments page',
    '  - See status: Pending / Approved / Rejected',
    '  - View test results and vaccination status',
    '  - Receive real-time notifications on status changes',
    '',
    'Step 5: Manage Profile',
    '  - Update personal information (name, phone, address, location)',
    '  - Delete account if needed',
]
add_bullet_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.5), flow_steps, font_size=14, color=DARK)

# ============================================================
# SLIDE 11: APPLICATION FLOW - HOSPITAL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Application Flow - Hospital Module')

flow_steps = [
    'Step 1: Register Account',
    '  - New hospitals register with name, email, password, address, location',
    '  - Initial status: "pending" - awaiting admin approval',
    '  - Admin receives notification of new registration',
    '',
    'Step 2: Login (after approval)',
    '  - Only approved hospitals can log in (checked in login controller)',
    '  - Rejected hospitals cannot access the system',
    '',
    'Step 3: Dashboard Overview',
    '  - View total and pending appointment requests',
    '  - See recent appointment activity',
    '',
    'Step 4: Manage Appointment Requests',
    '  - View all appointment requests on Requests page',
    '  - Approve: patient is notified',
    '  - Reject: patient is notified',
    '',
    'Step 5: Update Medical Records',
    '  - Update Test Result: Set as Negative or Positive',
    '  - Update Vaccination Status: Mark as Vaccinated',
    '  - Patient receives notifications for each update',
]
add_bullet_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.5), flow_steps, font_size=14, color=DARK)

# ============================================================
# SLIDE 12: APPLICATION FLOW - ADMIN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Application Flow - Admin Module')

flow_steps = [
    'Step 1: Login',
    '  - Admin logs in with email + password (web guard, role=admin)',
    '  - Pre-seeded admin: admin@covid.com / admin123',
    '',
    'Step 2: Dashboard',
    '  - Overview of system statistics:',
    '    - Total Patients registered',
    '    - Total Hospitals (approved/pending/rejected)',
    '    - Total Appointments booked',
    '    - Pending approvals count',
    '',
    'Step 3: Hospital Management',
    '  - View all registered hospitals with their status',
    '  - Approve hospital registrations',
    '  - Reject hospital registrations',
    '  - Hospitals are notified of their status changes',
    '',
    'Step 4: Vaccine Management',
    '  - View all vaccines in the system',
    '  - Add new vaccines',
    '  - Toggle vaccine status (Available / Unavailable)',
    '  - Approved hospitals notified when new vaccine added',
    '',
    'Step 5: Reports & Export',
    '  - View comprehensive appointment reports',
    '  - See patient & hospital details for each appointment',
    '  - Export all data as Excel (.xls) file',
]
add_bullet_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.5), flow_steps, font_size=14, color=DARK)

# ============================================================
# SLIDE 13: CONTROLLER ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Developer Guide - Controller Architecture')

# Auth controllers
add_textbox(slide, Inches(1), Inches(1.6), Inches(5), Inches(0.4), 'Auth Controllers:', font_size=20, bold=True, color=INDIGO)
auth_ctrls = [
    'PatientAuthController: login, register, logout',
    'HospitalAuthController: login (checks status), register, logout',
    'AdminAuthController: login, logout',
]
add_bullet_textbox(slide, Inches(1), Inches(2.1), Inches(5.5), Inches(2), auth_ctrls, font_size=14, color=DARK)

# Feature controllers
add_textbox(slide, Inches(7), Inches(1.6), Inches(5), Inches(0.4), 'Feature Controllers:', font_size=20, bold=True, color=INDIGO)
feat_ctrls = [
    'PatientController: dashboard, search, book, appointments, profile',
    'HospitalController: dashboard, requests, approve, reject, test-result, vaccination',
    'AdminController: dashboard, hospitals, vaccines, reports, export',
]
add_bullet_textbox(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(2), feat_ctrls, font_size=14, color=DARK)

# Middleware
add_textbox(slide, Inches(1), Inches(4.2), Inches(5), Inches(0.4), 'Middleware Guards:', font_size=20, bold=True, color=INDIGO)
mids = [
    'auth:web -> Protects patient & admin routes',
    'auth:hospital -> Protects hospital routes',
    'guest -> Login/Register pages (redirect if authenticated)',
]
add_bullet_textbox(slide, Inches(1), Inches(4.7), Inches(5.5), Inches(2), mids, font_size=14, color=DARK)

# Route structure
add_textbox(slide, Inches(7), Inches(4.2), Inches(5), Inches(0.4), 'Route Structure:', font_size=20, bold=True, color=INDIGO)
routes_info = [
    'patient.* -> /patient/* (13 routes)',
    'hospital.* -> /hospital/* (9 routes)',
    'admin.* -> /admin/* (11 routes)',
    'Welcome route -> / (landing page)',
]
add_bullet_textbox(slide, Inches(7), Inches(4.7), Inches(5.5), Inches(2), routes_info, font_size=14, color=DARK)

# ============================================================
# SLIDE 14: KEY CODE SNIPPETS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Developer Guide - Key Code Snippets')

add_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(0.4), 'Hospital Login Controller (status check):', font_size=18, bold=True, color=INDIGO)
code1 = """public function login(Request $request)
{
    $hospital = Hospital::where('email', $request->email)->first();
    if ($hospital && $hospital->status !== 'approved') {
        return back()->with('error', 'Account not approved yet.');
    }
    // ... proceed with authentication
}"""
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.1), Inches(11), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code1
p.font.size = Pt(11)
p.font.name = 'Consolas'
p.font.color.rgb = DARK

add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.4), 'Appointment Model (Relationships):', font_size=18, bold=True, color=INDIGO)
code2 = """class Appointment extends Model
{
    public function patient() {
        return $this->belongsTo(User::class, 'patient_id');
    }
    public function hospital() {
        return $this->belongsTo(Hospital::class);
    }
}"""
txBox = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code2
p.font.size = Pt(11)
p.font.name = 'Consolas'
p.font.color.rgb = DARK

add_textbox(slide, Inches(1), Inches(5.7), Inches(11), Inches(0.4), 'Notification System (Polymorphic):', font_size=18, bold=True, color=INDIGO)
code3 = """class Notification extends Model
{
    public function notifiable() {
        return $this->morphTo();
    }
    public function scopeUnread($query) {
        return $query->where('is_read', false);
    }
}"""
txBox = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code3
p.font.size = Pt(11)
p.font.name = 'Consolas'
p.font.color.rgb = DARK

# ============================================================
# SLIDE 15: KEY FEATURES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Key Features Highlight')

features = [
    'Multi-Role Authentication: Three separate auth systems with role-based access',
    'Hospital Approval Workflow: Admin must approve hospitals before they can operate',
    'Smart Search: Search approved hospitals by name or location',
    'Appointment Booking: Book COVID tests and vaccinations at verified hospitals',
    'Real-time Notifications: Polymorphic notification system for all user types',
    'Test Result Management: Hospitals can update test results (negative/positive)',
    'Vaccination Tracking: Track vaccination status per patient',
    'Vaccine Inventory: Admin manages vaccine availability',
    'Reporting & Export: Comprehensive reports with Excel export',
    'Responsive UI: Tailwind CSS with dark mode and theme customization',
    'Polymorphic Notifications: Single notifications table serving both User and Hospital models',
    'Secure Architecture: CSRF protection, bcrypt hashing, session management',
]
add_bullet_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.5), features, font_size=15, color=DARK)

# ============================================================
# SLIDE 16: USER GUIDE - PATIENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'User Guide - Patient')

user_steps = [
    '1. Visit the homepage and click "Get Started as Patient"',
    '2. Register with your name, email, password, phone, and address',
    '3. Log in using your email and password',
    '4. Navigate to "Search" from the navigation bar',
    '5. Search for hospitals by name or location',
    '6. Click "Book Appointment" on any hospital',
    '7. Select appointment type (Test / Vaccination) and choose a date',
    '8. Submit - your request is sent to the hospital',
    '9. View your appointments under "Appointments" menu',
    '10. Check status updates (Pending / Approved / Rejected)',
    '11. View test results and vaccination status once updated by hospital',
    '12. Edit your profile or delete your account from "Profile" page',
    '13. Click "Logout" when done',
]
add_bullet_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.5), user_steps, font_size=15, color=DARK)

# ============================================================
# SLIDE 17: USER GUIDE - HOSPITAL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'User Guide - Hospital')

user_steps = [
    '1. Visit the homepage and click "Register as Hospital"',
    '2. Fill in hospital name, email, password, address, and location',
    '3. Submit - your account status will be "Pending"',
    '4. Wait for admin approval (you will be notified)',
    '5. Once approved, log in with your email and password',
    '6. View your dashboard with appointment statistics',
    '7. Go to "Requests" to see all patient appointment requests',
    '8. Click "Approve" to confirm an appointment',
    '9. Click "Reject" to decline an appointment',
    '10. After appointment, update Test Result (Negative/Positive)',
    '11. Mark patient as "Vaccinated" after vaccination',
    '12. All status changes send notifications to patients',
    '13. Click "Logout" when done',
]
add_bullet_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.5), user_steps, font_size=15, color=DARK)

# ============================================================
# SLIDE 18: USER GUIDE - ADMIN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'User Guide - Admin')

user_steps = [
    '1. Visit homepage and click "Admin" button in navigation',
    '2. Log in with admin@covid.com / admin123',
    '3. Dashboard shows system-wide statistics',
    '4. Go to "Hospitals" to manage hospital registrations',
    '5. Approve or reject pending hospital accounts',
    '6. Go to "Vaccines" to manage vaccine inventory',
    '7. Add new vaccines by name',
    '8. Toggle vaccine availability (Available / Unavailable)',
    '9. Go to "Reports" to view all appointments in the system',
    '10. See patient names, hospital names, dates, and statuses',
    '11. Go to "Export" to download all data as Excel file',
    '12. Click "Logout" when done',
]
add_bullet_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.5), user_steps, font_size=15, color=DARK)

# ============================================================
# SLIDE 19: MODULE DESCRIPTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Developer Guide - Module Descriptions')

modules = [
    '1. Authentication Module:',
    '   - Handles registration & login for all 3 roles',
    '   - HospitalAuthController blocks unapproved hospitals',
    '   - AdminAuthController checks role=admin in users table',
    '',
    '2. Patient Module (PatientController):',
    '   - Dashboard, Hospital Search, Appointment Booking',
    '   - Appointment history & Profile management',
    '   - Account deletion with cascade',
    '',
    '3. Hospital Module (HospitalController):',
    '   - Dashboard with stats, Request management',
    '   - Appointment approval/rejection workflow',
    '   - Test result & vaccination status updates',
    '',
    '4. Admin Module (AdminController):',
    '   - System dashboard with aggregate statistics',
    '   - Hospital verification workflow',
    '   - Vaccine CRUD & toggle, Reports & Excel export',
    '',
    '5. Notification Module (Notification Model):',
    '   - Polymorphic notifications for User & Hospital',
    '   - Unread scope, markAsRead method',
    '   - Triggered on bookings, approvals, status changes',
]
add_bullet_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.5), modules, font_size=14, color=DARK)

# ============================================================
# SLIDE 20: TESTING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Testing Strategy')

testing = [
    'Test Framework: Laravel Test (PHPUnit wrapper)',
    'Database: SQLite in-memory for all tests',
    'Migration Strategy: DatabaseMigrations trait (fresh DB per test)',
    '',
    'Test Suites:',
    '  - Unit Tests: Individual component testing',
    '  - Feature Tests: Full request/response cycle testing',
    '',
    'Test Coverage Areas:',
    '  - Authentication flows (login, register, logout for all roles)',
    '  - Hospital status gating (unapproved hospitals cannot login)',
    '  - Appointment booking & approval workflow',
    '  - Profile update & account deletion',
    '  - Vaccine management (add, toggle)',
    '  - Report generation & data export',
    '',
    'Running Tests:',
    '  - composer test (runs config:clear + artisan test)',
    '  - php artisan test tests/Feature/ExampleTest.php',
    '  - php artisan test --filter=test_name',
]
add_bullet_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.5), testing, font_size=15, color=DARK)

# ============================================================
# SLIDE 21: SCREENSHOTS (placeholder)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, 'Screenshots')

# Screenshot placeholders
ss_data = [
    ('Welcome / Landing Page', Inches(0.8), Inches(1.6)),
    ('Patient Dashboard', Inches(4.8), Inches(1.6)),
    ('Hospital Requests', Inches(8.8), Inches(1.6)),
    ('Admin Reports', Inches(0.8), Inches(4.3)),
    ('Appointment Booking', Inches(4.8), Inches(4.3)),
    ('Admin Hospitals', Inches(8.8), Inches(4.3)),
]
for label, left, top in ss_data:
    shape = add_shape_bg(slide, left, top, Inches(3.6), Inches(2.3), LIGHT_BG)
    shape.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
    shape.line.width = Pt(1)
    add_textbox(slide, left, top + Inches(0.8), Inches(3.6), Inches(0.5), '[Screenshot]', font_size=14, bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(1.8), Inches(3.6), Inches(0.4), label, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5), 'Note: Please insert actual application screenshots into these placeholder areas.', font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 22: CONCLUSION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, INDIGO)
add_textbox(slide, Inches(1.5), Inches(2), Inches(10), Inches(1), 'Conclusion', font_size=44, bold=True, color=WHITE)
add_shape_bg(slide, Inches(1.5), Inches(3), Inches(2.5), Inches(0.06), WHITE)
conclusion = [
    'CovidCare successfully delivers a complete vaccination management ecosystem.',
    'The system streamlines the entire workflow from patient registration to vaccination tracking.',
    'Three-role architecture ensures proper access control and data privacy.',
    'Hospital approval workflow maintains quality and trust in the system.',
    'Real-time notifications keep all stakeholders informed of status changes.',
    'The modular Laravel architecture allows for easy extension and maintenance.',
]
add_bullet_textbox(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(3), conclusion, font_size=18, color=RGBColor(0xC7, 0xD2, 0xFE))

# ============================================================
# SLIDE 23: THANK YOU
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2), 'Thank You', font_size=52, bold=True, color=INDIGO, alignment=PP_ALIGN.CENTER)
add_shape_bg(slide, Inches(5.5), Inches(3.7), Inches(2.3), Inches(0.06), INDIGO)
add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.6), 'CovidCare - Vaccination Management System', font_size=22, color=DARK, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(5), Inches(10), Inches(0.5), 'Developed by: Muhammad Owais, Muhammad Hammad, Muhammad Hunain, Muhammad Hamayl, Muhammad Ali', font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.5), 'eProject - June 2026 to July 2026', font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# Save
output_path = 'D:\\laravel work\\V_S\\CovidCare_eProject_Report.pptx'
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
